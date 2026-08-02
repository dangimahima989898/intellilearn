from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File, Form, BackgroundTasks
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from database import get_db
from models.note import Note
from models.subject import Subject
from models.user import User
from models.uploaded_note import UploadedNote
from models.note_summary import NoteSummary
from models.summary_version import SummaryVersion
from models.summary_feedback import SummaryFeedback
from models.summary_review_log import SummaryReviewLog
from schemas.note import (
    NoteOut,
    UploadedNoteOut,
    NoteSummaryOut,
    SummaryVersionOut,
    SummaryFeedbackSubmit,
    SummaryApproveRequest,
    SummaryRejectRequest,
    NoteAnalyticsResponse,
    NoteAnalyticsItem
)
from utils.dependencies import get_current_user, require_admin, require_student, require_role
from utils.semester_filter import apply_semester_filter
from utils.llm_client import get_llm_response
from PyPDF2 import PdfReader
import uuid
import os
import shutil
from datetime import datetime
from utils.firebase import send_to_all_students

router = APIRouter()
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)
MAX_SIZE = int(os.getenv("MAX_FILE_SIZE_MB", "10")) * 1024 * 1024
ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.ppt', '.pptx'}

def secure_filename(filename: str) -> str:
    return "".join([c for c in filename if c.isalpha() or c.isdigit() or c in (' ', '.', '_', '-')]).rstrip()

def serialize_note(note: Note, subject_name: str, uploaded_by_name: str) -> dict:
    return {
        "id": note.id,
        "title": note.title,
        "subject_id": note.subject_id,
        "subject_name": subject_name,
        "file_url": note.file_url,
        "file_type": note.file_type,
        "file_size_kb": note.file_size_kb,
        "uploaded_by_name": uploaded_by_name,
        "created_at": note.created_at,
        "summary": note.summary,
        "download_count": note.download_count,
        "course_id": note.course_id,
        "semester_number": note.semester_number,
    }

@router.post("/upload", response_model=NoteOut, status_code=status.HTTP_201_CREATED)
async def upload_note(
    title: str = Form(...),
    subject_id: uuid.UUID = Form(...),
    file: UploadFile = File(...),
    course_id: uuid.UUID = Form(None),
    semester_number: int = Form(None),
    db: Session = Depends(get_db),
    current_user = Depends(require_role("super_admin", "hod", "faculty"))
):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}")
        
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    
    if file_size > MAX_SIZE:
        raise HTTPException(status_code=400, detail=f"File too large. Max size is {MAX_SIZE / (1024*1024)}MB")
        
    subject = db.query(Subject).filter(Subject.id == subject_id).first()
    if not subject:
        raise HTTPException(status_code=404, detail="Subject not found")
        
    # Faculty ownership validation
    if current_user.role == "faculty":
        from models.faculty_subject_assignment import FacultySubjectAssignment
        assignment = db.query(FacultySubjectAssignment).filter(
            FacultySubjectAssignment.faculty_id == current_user.id,
            FacultySubjectAssignment.subject_id == subject_id
        ).first()
        if not assignment:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied. You are not assigned to this subject."
            )

    # Enforce course_id and semester_number match subject metadata
    if course_id and course_id != subject.course_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Cross-semester/cross-course uploads are not allowed. Selected course does not match subject."
        )
    if semester_number and semester_number != subject.semester_number:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Cross-semester/cross-course uploads are not allowed. Selected semester does not match subject."
        )
    
    course_id = subject.course_id
    semester_number = subject.semester_number
    safe_name = secure_filename(file.filename)
    unique_filename = f"{uuid.uuid4()}_{safe_name}"
    file_path = os.path.join(UPLOAD_DIR, unique_filename)
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    new_note = Note(
        id=uuid.uuid4(),
        subject_id=subject_id,
        uploaded_by=current_user.id,
        title=title,
        file_url=f"/uploads/{unique_filename}",
        file_type=ext[1:].lower(),
        file_size_kb=file_size // 1024,
        download_count=0,
        course_id=course_id,
        semester_number=semester_number
    )
    db.add(new_note)
    db.commit()
    db.refresh(new_note)
    
    send_to_all_students(
        db=db,
        title="📚 New Study Material Available",
        body=f"New notes for {subject.name}: {new_note.title}",
        data={"note_id": str(new_note.id), "type": "note"}
    )
    
    return serialize_note(new_note, subject.name, current_user.name)

@router.get("", response_model=list[NoteOut])
def get_notes(
    subject_id: uuid.UUID = None,
    course_id: uuid.UUID = None,
    semester: int = None,
    include_archived: bool = False,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    query = db.query(Note, Subject.name.label("subject_name"), User.name.label("uploaded_by_name"))\
              .join(Subject, Note.subject_id == Subject.id)\
              .join(User, Note.uploaded_by == User.id)
              
    if not include_archived:
        query = query.filter(Subject.is_archived == False)
              
    if current_user.role == "student":
        if subject_id:
            # When fetching notes for a specific subject, just filter by subject.
            # The subject is already verified to belong to the student's semester
            # via the /subjects endpoint, so no extra course/semester filter needed.
            query = query.filter(Note.subject_id == subject_id)
        else:
            # Browsing all notes — apply semester filter so students only see their semester
            query = apply_semester_filter(query, Note, current_user)
    else:
        if course_id:
            query = query.filter(Note.course_id == course_id)
        if semester:
            query = query.filter(Note.semester_number == semester)
        if subject_id:
            query = query.filter(Note.subject_id == subject_id)
        
    results = query.all()
    
    return [
        serialize_note(note, subject_name, uploaded_by_name)
        for note, subject_name, uploaded_by_name in results
    ]


@router.get("/draft-summaries")
def get_draft_summaries(
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    results = db.query(NoteSummary)\
        .join(UploadedNote, NoteSummary.note_id == UploadedNote.id)\
        .filter(NoteSummary.status.in_(["DRAFT", "EXTRACTING", "SUMMARIZING", "FAILED", "REJECTED"]))\
        .order_by(NoteSummary.created_at.desc()).all()
        
    response = []
    for summary in results:
        latest_version = db.query(SummaryVersion)\
            .filter(SummaryVersion.summary_id == summary.id)\
            .order_by(SummaryVersion.version_number.desc()).first()
            
        summary_text = latest_version.summary_text if latest_version else ""
        
        response.append({
            "id": str(summary.id),
            "note_id": str(summary.note_id),
            "title": summary.note.title,
            "subject_id": str(summary.note.subject_id),
            "subject_name": summary.note.subject.name if summary.note.subject else "",
            "unit": summary.note.unit,
            "file_url": summary.note.file_url,
            "status": summary.status,
            "current_version": summary.current_version,
            "summary_text": summary_text,
            "rejection_comment": summary.rejection_comment,
            "created_at": summary.created_at,
            "updated_at": summary.updated_at
        })
    return response


@router.get("/{id}", response_model=NoteOut)
def get_note(id: uuid.UUID, db: Session = Depends(get_db), current_user = Depends(get_current_user)):
    result = db.query(Note, Subject.name.label("subject_name"), User.name.label("uploaded_by_name"))\
               .join(Subject, Note.subject_id == Subject.id)\
               .join(User, Note.uploaded_by == User.id)\
               .filter(Note.id == id, Subject.is_archived == False).first()
               
    if not result:
        raise HTTPException(status_code=404, detail="Note not found")
        
    note, subject_name, uploaded_by_name = result
    return serialize_note(note, subject_name, uploaded_by_name)

@router.get("/download/{id}")
def download_note(id: uuid.UUID, db: Session = Depends(get_db), current_user = Depends(get_current_user)):
    note = db.query(Note).join(Subject).filter(Note.id == id, Subject.is_archived == False).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
        
    note.download_count += 1
    db.commit()
    
    file_path = os.path.join(".", note.file_url.lstrip("/"))
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on server")
        
    return FileResponse(path=file_path, filename=os.path.basename(file_path))

@router.delete("/{id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_note(id: uuid.UUID, db: Session = Depends(get_db), current_user = Depends(require_role("super_admin", "hod", "faculty"))):
    note = db.query(Note).filter(Note.id == id).first()
    if note:
        file_path = os.path.join(".", note.file_url.lstrip("/"))
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as e:
                print(f"Error removing file {file_path}: {e}")
        db.delete(note)
        db.commit()
        return None

    uploaded_note = db.query(UploadedNote).filter(UploadedNote.id == id).first()
    if uploaded_note:
        file_path = os.path.join(".", uploaded_note.file_url.lstrip("/"))
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as e:
                print(f"Error removing file {file_path}: {e}")
        db.delete(uploaded_note)
        db.commit()
        return None

    raise HTTPException(status_code=404, detail="Note not found")

@router.post("/{id}/summarize")
async def summarize_note(
    id: uuid.UUID,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    note = db.query(Note).filter(Note.id == id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")
        
    file_path = os.path.join(".", note.file_url.lstrip("/"))
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found on server")
        
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            text = " ".join([page.extract_text() or "" for page in reader.pages])
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext in ['.ppt', '.pptx']:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type for summarization: {ext}")
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return {"summary": "Could not extract text from this file. The file may be scanned/image-based.", "note_id": str(id)}
        
    text = text.strip()
    if not text:
        return {"summary": "Could not extract text from this file. The file may be scanned/image-based.", "note_id": str(id)}
        
    # Truncate text to first 6000 chars for safety
    truncated_text = text[:6000]
    
    prompt = f"""You are an expert academic summarizer for MCA students. Summarize the following study material concisely.

Format your summary as:
## Key Concepts
- [bullet point 1]
- [bullet point 2]
...

## Important Definitions
- Term: Definition
...

## Key Points to Remember
1. Point 1
2. Point 2
...

Keep the summary under 300 words total. Focus on what students need to know for exams.

Study Material:
{truncated_text}"""

    try:
        summary_text = await get_llm_response(
            messages=[{"role": "user", "content": "Please generate the summary now."}],
            system_prompt=prompt,
            max_tokens=800
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=f"LLM summarization failed: {str(e)}"
        )
        
    note.summary = summary_text.strip()
    db.commit()
    db.refresh(note)
    
    return {
        "summary": note.summary,
        "note_id": str(note.id),
        "generated_at": datetime.utcnow()
    }


# ── Smart Notes Summarizer Endpoints ──────────────────────────────────────────

async def process_note_background(note_id: uuid.UUID, file_path: str, subject_name: str, unit: str, session_factory):
    db = session_factory()
    try:
        from utils.pdf_processor import extract_with_ocr_fallback, clean_extracted_text
        from utils.summarizer_service import generate_summary
        
        # 1. Update status to EXTRACTING
        summary = db.query(NoteSummary).filter(NoteSummary.note_id == note_id).first()
        if summary:
            summary.status = "EXTRACTING"
            db.commit()
            
        # Extract text
        raw_text = extract_with_ocr_fallback(file_path)
        cleaned_text = clean_extracted_text(raw_text)
        
        # Save raw_text
        note = db.query(UploadedNote).filter(UploadedNote.id == note_id).first()
        if note:
            note.raw_text = cleaned_text
            note.file_size_kb = os.path.getsize(file_path) // 1024
            db.commit()
            
        # Update status to SUMMARIZING
        summary = db.query(NoteSummary).filter(NoteSummary.note_id == note_id).first()
        if summary:
            summary.status = "SUMMARIZING"
            db.commit()
            
        # 2. Summarize
        summary_text = await generate_summary(
            text=cleaned_text or "No readable text content found in PDF.",
            subject=subject_name,
            unit=unit
        )
        
        # Save version 1
        version = SummaryVersion(
            id=uuid.uuid4(),
            summary_id=summary.id,
            version_number=1,
            summary_text=summary_text,
            created_by_ai=True
        )
        db.add(version)
        
        # Update status to DRAFT
        summary.status = "DRAFT"
        db.commit()

        # Trigger RAG pipeline after raw_text is saved and summary is generated
        if note and note.raw_text and note.raw_text.strip():
            from utils.semantic_chunker import SemanticChunker
            from utils.embedding_service import EmbeddingService
            from utils.processing_pipeline import ProcessingPipeline

            chunker = SemanticChunker()
            embedding_svc = EmbeddingService()
            pipeline = ProcessingPipeline(chunker, embedding_svc)

            try:
                await pipeline.process_note(note.id, note.raw_text, note.subject_id)
            except Exception as rag_err:
                print(f"[RAG Pipeline] Error processing note {note_id}: {rag_err}")
    except Exception as e:
        print(f"Error processing note background for note {note_id}: {e}")
        import traceback
        traceback.print_exc()
        summary = db.query(NoteSummary).filter(NoteSummary.note_id == note_id).first()
        if summary:
            summary.status = "FAILED"
            db.commit()
    finally:
        db.close()


@router.post("/summary-upload", status_code=status.HTTP_201_CREATED)
async def upload_summary_note(
    background_tasks: BackgroundTasks,
    subject_id: uuid.UUID = Form(...),
    unit: str = Form(...),
    files: list[UploadFile] = File(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    subject = db.query(Subject).filter(Subject.id == subject_id).first()
    if not subject:
        raise HTTPException(status_code=404, detail="Subject not found")

    # Faculty ownership validation
    if current_user.role == "faculty":
        from models.faculty_subject_assignment import FacultySubjectAssignment
        assignment = db.query(FacultySubjectAssignment).filter(
            FacultySubjectAssignment.faculty_id == current_user.id,
            FacultySubjectAssignment.subject_id == subject_id
        ).first()
        if not assignment:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied. You are not assigned to this subject."
            )

    response_notes = []
    
    for file in files:
        ext = os.path.splitext(file.filename)[1].lower()
        if ext != '.pdf':
            raise HTTPException(status_code=400, detail="Only PDF files are supported for notes summarizer.")
            
        safe_name = secure_filename(file.filename)
        unique_filename = f"{uuid.uuid4()}_{safe_name}"
        file_path = os.path.join(UPLOAD_DIR, unique_filename)
        
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
            
        new_note = UploadedNote(
            id=uuid.uuid4(),
            title=os.path.splitext(safe_name)[0],
            subject_id=subject_id,
            unit=unit,
            file_url=f"/uploads/{unique_filename}",
            uploaded_by=current_user.id
        )
        db.add(new_note)
        db.commit()
        db.refresh(new_note)
        
        new_summary = NoteSummary(
            id=uuid.uuid4(),
            note_id=new_note.id,
            status="EXTRACTING",
            current_version=1
        )
        db.add(new_summary)
        db.commit()
        db.refresh(new_summary)
        
        from database import SessionLocal
        background_tasks.add_task(
            process_note_background,
            new_note.id,
            file_path,
            subject.name,
            unit,
            SessionLocal
        )
        
        response_notes.append({
            "id": str(new_note.id),
            "title": new_note.title,
            "subject_id": str(new_note.subject_id),
            "subject_name": subject.name,
            "unit": new_note.unit,
            "file_url": new_note.file_url,
            "created_at": new_note.created_at
        })
        
    return response_notes



@router.post("/approve/{summary_id}")
def approve_summary(
    summary_id: uuid.UUID,
    payload: SummaryApproveRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    summary = db.query(NoteSummary).filter(NoteSummary.id == summary_id).first()
    if not summary:
        raise HTTPException(status_code=404, detail="Summary not found")
        
    # Faculty ownership validation
    if current_user.role == "faculty":
        from models.faculty_subject_assignment import FacultySubjectAssignment
        assignment = db.query(FacultySubjectAssignment).filter(
            FacultySubjectAssignment.faculty_id == current_user.id,
            FacultySubjectAssignment.subject_id == summary.note.subject_id
        ).first()
        if not assignment:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied. You are not assigned to this subject."
            )

    latest_version = db.query(SummaryVersion)\
        .filter(SummaryVersion.summary_id == summary_id)\
        .order_by(SummaryVersion.version_number.desc()).first()
        
    is_edited = latest_version is None or latest_version.summary_text.strip() != payload.summary_text.strip()
    action = "EDITED & APPROVED" if is_edited else "APPROVED"
    
    if is_edited:
        summary.current_version += 1
        new_version = SummaryVersion(
            id=uuid.uuid4(),
            summary_id=summary.id,
            version_number=summary.current_version,
            summary_text=payload.summary_text,
            created_by_ai=False,
            approved_by=current_user.id,
            approved_at=datetime.utcnow()
        )
        db.add(new_version)
    else:
        if latest_version:
            latest_version.approved_by = current_user.id
            latest_version.approved_at = datetime.utcnow()
            db.add(latest_version)
            
    summary.status = "APPROVED"
    db.add(summary)
    
    log = SummaryReviewLog(
        id=uuid.uuid4(),
        summary_id=summary.id,
        faculty_id=current_user.id,
        action=action,
        comment=None
    )
    db.add(log)
    db.commit()
    db.refresh(summary)
    
    try:
        send_to_all_students(
            db=db,
            title="📚 Approved Summary Available",
            body=f"New summary approved for {summary.note.subject.name} - {summary.note.unit}: {summary.note.title}",
            data={"summary_id": str(summary.id), "type": "summary"}
        )
    except Exception as e:
        print(f"Failed to send notification: {e}")
        
    return {"message": "Summary approved successfully", "status": summary.status}


@router.post("/reject/{summary_id}")
def reject_summary(
    summary_id: uuid.UUID,
    payload: SummaryRejectRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    summary = db.query(NoteSummary).filter(NoteSummary.id == summary_id).first()
    if not summary:
        raise HTTPException(status_code=404, detail="Summary not found")
        
    # Faculty ownership validation
    if current_user.role == "faculty":
        from models.faculty_subject_assignment import FacultySubjectAssignment
        assignment = db.query(FacultySubjectAssignment).filter(
            FacultySubjectAssignment.faculty_id == current_user.id,
            FacultySubjectAssignment.subject_id == summary.note.subject_id
        ).first()
        if not assignment:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Access denied. You are not assigned to this subject."
            )

    summary.status = "REJECTED"
    summary.rejection_comment = payload.reason
    db.add(summary)
    
    log = SummaryReviewLog(
        id=uuid.uuid4(),
        summary_id=summary.id,
        faculty_id=current_user.id,
        action="REJECTED",
        comment=payload.reason
    )
    db.add(log)
    db.commit()
    db.refresh(summary)
    
    return {"message": "Summary rejected successfully", "status": summary.status}


@router.get("/summaries/{subject_id}")
def get_approved_summaries(
    subject_id: uuid.UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    results = db.query(NoteSummary)\
        .join(UploadedNote, NoteSummary.note_id == UploadedNote.id)\
        .filter(UploadedNote.subject_id == subject_id, NoteSummary.status == "APPROVED")\
        .order_by(NoteSummary.updated_at.desc()).all()
        
    response = []
    for summary in results:
        latest_version = db.query(SummaryVersion)\
            .filter(SummaryVersion.summary_id == summary.id, SummaryVersion.approved_by != None)\
            .order_by(SummaryVersion.version_number.desc()).first()
            
        if not latest_version:
            latest_version = db.query(SummaryVersion)\
                .filter(SummaryVersion.summary_id == summary.id)\
                .order_by(SummaryVersion.version_number.desc()).first()
                
        summary_text = latest_version.summary_text if latest_version else ""
        approver_name = latest_version.approver.name if (latest_version and latest_version.approver) else "Faculty"
        approved_at = latest_version.approved_at if latest_version else None
        
        response.append({
            "id": str(summary.id),
            "note_id": str(summary.note_id),
            "title": summary.note.title,
            "unit": summary.note.unit,
            "file_url": summary.note.file_url,
            "current_version": summary.current_version,
            "summary_text": summary_text,
            "views_count": summary.views_count,
            "avg_read_time_seconds": summary.avg_read_time_seconds,
            "helpful_count": summary.helpful_count,
            "not_helpful_count": summary.not_helpful_count,
            "approver_name": approver_name,
            "approved_at": approved_at,
            "created_at": summary.created_at
        })
    return response


@router.post("/summaries/view/{summary_id}")
def increment_summary_views(
    summary_id: uuid.UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    summary = db.query(NoteSummary).filter(NoteSummary.id == summary_id).first()
    if not summary:
        raise HTTPException(status_code=404, detail="Summary not found")
    summary.views_count += 1
    db.commit()
    return {"views_count": summary.views_count}


@router.post("/feedback/{summary_id}")
def submit_summary_feedback(
    summary_id: uuid.UUID,
    payload: SummaryFeedbackSubmit,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_student)
):
    summary = db.query(NoteSummary).filter(NoteSummary.id == summary_id).first()
    if not summary:
        raise HTTPException(status_code=404, detail="Summary not found")
        
    feedback = db.query(SummaryFeedback)\
        .filter(SummaryFeedback.summary_id == summary_id, SummaryFeedback.student_id == current_user.id).first()
        
    if feedback:
        feedback.is_helpful = payload.is_helpful
        feedback.time_spent_seconds = payload.time_spent_seconds
        feedback.created_at = datetime.utcnow()
    else:
        feedback = SummaryFeedback(
            id=uuid.uuid4(),
            summary_id=summary_id,
            student_id=current_user.id,
            is_helpful=payload.is_helpful,
            time_spent_seconds=payload.time_spent_seconds
        )
        db.add(feedback)
        
    db.commit()
    
    all_feedbacks = db.query(SummaryFeedback).filter(SummaryFeedback.summary_id == summary_id).all()
    helpful = sum(1 for f in all_feedbacks if f.is_helpful)
    not_helpful = sum(1 for f in all_feedbacks if not f.is_helpful)
    total_time = sum(f.time_spent_seconds for f in all_feedbacks)
    avg_time = total_time / len(all_feedbacks) if all_feedbacks else 0.0
    
    summary.helpful_count = helpful
    summary.not_helpful_count = not_helpful
    summary.total_read_time_seconds = float(total_time)
    summary.avg_read_time_seconds = float(avg_time)
    
    db.commit()
    db.refresh(summary)
    
    return {
        "helpful_count": summary.helpful_count,
        "not_helpful_count": summary.not_helpful_count,
        "avg_read_time_seconds": summary.avg_read_time_seconds
    }


@router.get("/versions/{summary_id}")
def get_summary_versions(
    summary_id: uuid.UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    versions = db.query(SummaryVersion)\
        .filter(SummaryVersion.summary_id == summary_id)\
        .order_by(SummaryVersion.version_number.desc()).all()
        
    response = []
    for v in versions:
        response.append({
            "id": str(v.id),
            "version_number": v.version_number,
            "summary_text": v.summary_text,
            "created_by_ai": v.created_by_ai,
            "approved_by_name": v.approver.name if v.approver else None,
            "approved_at": v.approved_at,
            "created_at": v.created_at
        })
    return response


@router.post("/regenerate/{note_id}")
async def force_regenerate_summary(
    note_id: uuid.UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    try:
        summary = await regenerate_summary(db, note_id, created_by_ai=True)
        latest_version = db.query(SummaryVersion)\
            .filter(SummaryVersion.summary_id == summary.id)\
            .order_by(SummaryVersion.version_number.desc()).first()
            
        return {
            "message": "Summary regenerated successfully",
            "summary_id": str(summary.id),
            "status": summary.status,
            "current_version": summary.current_version,
            "summary_text": latest_version.summary_text if latest_version else ""
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to regenerate summary: {str(e)}")


@router.get("/download-summary-pdf/{summary_id}")
def download_summary_pdf_file(
    summary_id: uuid.UUID,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    summary = db.query(NoteSummary).filter(NoteSummary.id == summary_id).first()
    if not summary:
        raise HTTPException(status_code=404, detail="Summary not found")
        
    latest_version = db.query(SummaryVersion)\
        .filter(SummaryVersion.summary_id == summary_id)\
        .order_by(SummaryVersion.version_number.desc()).first()
        
    if not latest_version:
        raise HTTPException(status_code=404, detail="No summary version content found")
        
    subject_code = summary.note.subject.code if summary.note.subject else "MCA"
    subject_name = summary.note.subject.name if summary.note.subject else "Subject"
    approver_name = latest_version.approver.name if latest_version.approver else "Assigned Professor"
    approval_date = latest_version.approved_at.strftime("%Y-%m-%d") if latest_version.approved_at else datetime.utcnow().strftime("%Y-%m-%d")
    
    metadata = {
        "title": summary.note.title,
        "subject_name": subject_name,
        "subject_code": subject_code,
        "unit": summary.note.unit,
        "approver_name": approver_name,
        "approval_date": approval_date
    }
    
    pdf_filename = f"summary_{summary_id}_{summary.current_version}.pdf"
    pdf_path = os.path.join(UPLOAD_DIR, pdf_filename)
    
    try:
        generate_summary_pdf(latest_version.summary_text, metadata, pdf_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PDF: {str(e)}")
        
    if not os.path.exists(pdf_path):
        raise HTTPException(status_code=404, detail="Generated PDF not found on server")
        
    return FileResponse(path=pdf_path, filename=f"{summary.note.title}_summary.pdf", media_type="application/pdf")


@router.get("/admin/analytics", response_model=NoteAnalyticsResponse)
def get_notes_summary_analytics(
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    summaries = db.query(NoteSummary).all()
    
    total_uploaded = db.query(UploadedNote).count()
    total_approved = db.query(NoteSummary).filter(NoteSummary.status == "APPROVED").count()
    total_pending = db.query(NoteSummary).filter(NoteSummary.status.in_(["DRAFT", "EXTRACTING", "SUMMARIZING"])).count()
    total_rejected = db.query(NoteSummary).filter(NoteSummary.status == "REJECTED").count()
    
    items = []
    for summary in summaries:
        total_votes = summary.helpful_count + summary.not_helpful_count
        helpfulness = (summary.helpful_count / total_votes * 100.0) if total_votes > 0 else 100.0
        
        items.append(NoteAnalyticsItem(
            note_id=summary.note_id,
            title=summary.note.title,
            subject_name=summary.note.subject.name if summary.note.subject else "Unknown",
            unit=summary.note.unit,
            views_count=summary.views_count,
            avg_read_time_seconds=summary.avg_read_time_seconds,
            helpfulness_percentage=helpfulness,
            status=summary.status,
            helpful_count=summary.helpful_count,
            not_helpful_count=summary.not_helpful_count
        ))
        
    return NoteAnalyticsResponse(
        total_uploaded=total_uploaded,
        total_approved=total_approved,
        total_pending=total_pending,
        total_rejected=total_rejected,
        items=items
    )


# ── Note Update Endpoint (triggers RAG reprocessing) ──────────────────────────

@router.put("/update-text/{note_id}")
async def update_note_raw_text(
    note_id: uuid.UUID,
    background_tasks: BackgroundTasks,
    raw_text: str = Form(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(require_role("super_admin", "hod", "faculty"))
):
    """Update a note's raw_text and trigger RAG pipeline reprocessing if text changed."""
    note = db.query(UploadedNote).filter(UploadedNote.id == note_id).first()
    if not note:
        raise HTTPException(status_code=404, detail="Note not found")

    old_raw_text = note.raw_text
    note.raw_text = raw_text
    db.commit()
    db.refresh(note)

    # Trigger RAG reprocessing if raw_text changed and has content
    if raw_text.strip() and raw_text != old_raw_text:
        background_tasks.add_task(
            _reprocess_note_background,
            note.id,
            raw_text,
            note.subject_id
        )

    return {"message": "Note text updated", "note_id": str(note.id)}


async def _reprocess_note_background(note_id: uuid.UUID, raw_text: str, subject_id: uuid.UUID):
    """Background task to reprocess a note through the RAG pipeline."""
    from utils.semantic_chunker import SemanticChunker
    from utils.embedding_service import EmbeddingService
    from utils.processing_pipeline import ProcessingPipeline

    chunker = SemanticChunker()
    embedding_svc = EmbeddingService()
    pipeline = ProcessingPipeline(chunker, embedding_svc)

    try:
        await pipeline.reprocess_note(note_id, raw_text, subject_id)
    except Exception as e:
        print(f"[RAG Pipeline] Error reprocessing note {note_id}: {e}")
